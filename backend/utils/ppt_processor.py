from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
import os
import uuid
from copy import deepcopy
import platform
import shutil
import subprocess
try:
    import spacy  # type: ignore
except Exception:
    spacy = None  # type: ignore
import requests
import json
import re
from PIL import Image
import io

# Optional Windows-only COM automation (guarded)
WIN32_AVAILABLE = False
if platform.system() == "Windows":
    try:
        import win32com.client  # type: ignore
        import pythoncom  # type: ignore
        WIN32_AVAILABLE = True
    except Exception:
        WIN32_AVAILABLE = False

# Load spaCy model for text analysis (graceful fallback if unavailable)
if spacy is not None:
    try:
        nlp = spacy.load("en_core_web_sm")
    except Exception:
        try:
            from spacy.cli import download  # type: ignore
            download("en_core_web_sm")
            nlp = spacy.load("en_core_web_sm")
        except Exception:
            nlp = None
else:
    nlp = None

UNSPLASH_API_KEY = os.getenv("UNSPLASH_API_KEY")
PEXELS_API_KEY = os.getenv("PEXELS_API_KEY")

def fetch_image_from_pexels(query, save_dir="Uploads/images"):
    """Fallback provider using Pexels API."""
    try:
        api_key = (os.getenv("PEXELS_API_KEY") or PEXELS_API_KEY or "").strip()
        if not api_key:
            return None
        q = str(query or "").strip()
        q = re.sub(r"[\n\r]+", " ", q)
        q = re.sub(r"[\"'(){}\[\]:]", "", q)
        q = re.sub(r"\s+", " ", q)
        q = q[:120] or "abstract background"
        url = "https://api.pexels.com/v1/search"
        headers = {"Authorization": api_key}
        params = {"query": q, "per_page": 1}
        resp = requests.get(url, headers=headers, params=params, timeout=15)
        try:
            resp.raise_for_status()
        except Exception:
            print(f"Pexels API error ({resp.status_code}) for query '{q}': {resp.text[:180]}")
            return None
        photos = resp.json().get("photos", [])
        if not photos:
            return None
        image_url = photos[0].get("src", {}).get("large") or photos[0].get("src", {}).get("original")
        if not image_url:
            return None
        image_filename = f"slideimg_{uuid.uuid4().hex[:8]}.jpg"
        os.makedirs(save_dir, exist_ok=True)
        image_path = os.path.join(save_dir, image_filename)
        img = requests.get(image_url, timeout=20)
        img.raise_for_status()
        with open(image_path, "wb") as f:
            f.write(img.content)
        print(f"Successfully fetched image from Pexels: {image_path}")
        return image_path
    except Exception as e:
        print(f"Pexels fetch failed for '{query}': {e}")
        return None

def fetch_image_from_unsplash(query, save_dir="Uploads/images"):
    try:
        # Re-read API key each call in case env changed
        api_key = (os.getenv("UNSPLASH_API_KEY") or UNSPLASH_API_KEY or "").strip()
        if not api_key:
            print("No Unsplash API key found, using placeholder")
            return create_placeholder_image(query, save_dir)

        # Sanitize and trim verbose prompts into keywords
        q = str(query or "").strip()
        q = re.sub(r"[\n\r]+", " ", q)
        q = re.sub(r"[\"'(){}\[\]:]", "", q)
        q = re.sub(r"\s+", " ", q)
        q = q[:120]
        if not q:
            q = "abstract background"

        url = "https://api.unsplash.com/search/photos"
        headers = {
            "Authorization": f"Client-ID {api_key}",
            "Accept-Version": "v1",
        }
        params = {"query": q, "per_page": 1, "content_filter": "high"}
        response = requests.get(url, headers=headers, params=params, timeout=15)
        try:
            response.raise_for_status()
        except Exception as e:
            # Provide clearer diagnostics for 401/403
            print(f"Unsplash API error ({response.status_code}) for query '{q}': {response.text[:180]}")
            raise

        results = response.json().get("results", [])
        if results:
            image_url = results[0]["urls"].get("regular") or results[0]["urls"].get("full")
            image_filename = f"slideimg_{uuid.uuid4().hex[:8]}.jpg"
            os.makedirs(save_dir, exist_ok=True)
            image_path = os.path.join(save_dir, image_filename)
            img = requests.get(image_url, timeout=20)
            img.raise_for_status()
            with open(image_path, "wb") as f:
                f.write(img.content)
            print(f"Successfully fetched image from Unsplash: {image_path}")
            return image_path
        else:
            print(f"No images found for '{q}' on Unsplash, trying Pexels...")
            p = fetch_image_from_pexels(q, save_dir)
            if p:
                return p
            return create_placeholder_image(q, save_dir)
    except Exception as e:
        print(f"Image fetch failed for '{query}': {e}. Trying Pexels...")
        p = fetch_image_from_pexels(query, save_dir)
        if p:
            return p
        return create_placeholder_image(query, save_dir)

def create_placeholder_image(query, save_dir="Uploads/images"):
    """
    Create a beautiful placeholder image with gradient background and text
    """
    try:
        from PIL import Image, ImageDraw, ImageFont
        
        # Create a beautiful gradient placeholder image
        width, height = 800, 600
        
        # Create gradient background
        img = Image.new('RGB', (width, height), color=(41, 128, 185))  # Blue gradient start
        draw = ImageDraw.Draw(img)
        
        # Create gradient effect
        for y in range(height):
            r = int(41 + (y / height) * 30)  # Blue to lighter blue
            g = int(128 + (y / height) * 40)
            b = int(185 + (y / height) * 20)
            draw.line([(0, y), (width, y)], fill=(r, g, b))
        
        # Add a subtle pattern
        for i in range(0, width, 50):
            for j in range(0, height, 50):
                draw.ellipse([i, j, i+2, j+2], fill=(255, 255, 255, 30))
        
        # Add border
        draw.rectangle([0, 0, width-1, height-1], outline=(255, 255, 255), width=5)
        
        # Add icon placeholder
        icon_size = 80
        icon_x = (width - icon_size) // 2
        icon_y = (height - icon_size) // 2 - 50
        draw.ellipse([icon_x, icon_y, icon_x + icon_size, icon_y + icon_size], 
                    outline=(255, 255, 255), width=3, fill=(255, 255, 255, 50))
        
        # Add text
        try:
            # Try to use a default font
            font = ImageFont.load_default()
        except:
            font = None
        
        # Draw the query text
        text = query[:40] + "..." if len(query) > 40 else query
        text_bbox = draw.textbbox((0, 0), text, font=font)
        text_width = text_bbox[2] - text_bbox[0]
        text_height = text_bbox[3] - text_bbox[1]
        
        x = (width - text_width) // 2
        y = (height - text_height) // 2 + 50
        
        # Add text shadow
        draw.text((x+2, y+2), text, fill=(0, 0, 0, 100), font=font)
        draw.text((x, y), text, fill=(255, 255, 255), font=font)
        
        # Add "Image Placeholder" text
        placeholder_text = "Image Placeholder"
        placeholder_bbox = draw.textbbox((0, 0), placeholder_text, font=font)
        placeholder_width = placeholder_bbox[2] - placeholder_bbox[0]
        placeholder_x = (width - placeholder_width) // 2
        placeholder_y = y + text_height + 20
        
        draw.text((placeholder_x, placeholder_y), placeholder_text, fill=(255, 255, 255, 150), font=font)
        
        # Save the image
        image_filename = f"placeholder_{uuid.uuid4().hex[:8]}.png"
        os.makedirs(save_dir, exist_ok=True)
        image_path = os.path.join(save_dir, image_filename)
        img.save(image_path)
        
        print(f"Created beautiful placeholder image: {image_path}")
        return image_path
    except Exception as e:
        print(f"Failed to create placeholder image: {e}")
        return None

def add_custom_slide(prs, slide_data, image_path=None):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    layout = slide_data.get("layout", "").lower()
    
    # Add image first if available
    if layout == "image left, text right" and image_path:
        slide.shapes.add_picture(image_path, Inches(0.2), Inches(1), Inches(4), Inches(4))
        txBox = slide.shapes.add_textbox(Inches(4.5), Inches(1), Inches(5), Inches(4))
    elif layout == "image right, text left" and image_path:
        txBox = slide.shapes.add_textbox(Inches(0.2), Inches(1), Inches(5), Inches(4))
        slide.shapes.add_picture(image_path, Inches(5.5), Inches(1), Inches(4), Inches(4))
    elif layout == "full image" and image_path:
        slide.shapes.add_picture(image_path, Inches(0.2), Inches(0.2), Inches(9), Inches(6))
        txBox = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1.2))
    else:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
    
    tf = txBox.text_frame
    tf.text = slide_data.get("title", "")
    for bullet in slide_data.get("bullets", []):
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1
    
    # Detect contrast and apply appropriate text color
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    contrast_type = get_contrast_text_color(slide, slide_width, slide_height)
    text_color = get_text_color_for_contrast(contrast_type, slide_data.get("theme", ["#003087", "#FFFFFF"]))
    
    # Apply contrast-based text color to all paragraphs
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = text_color
    
    return slide

def add_two_column_content(slide, slide_data, presentation, slide_width, slide_height):
    """Render a two-column layout with a single centered bold title at the top,
    then two body text columns (left/right) below. Sizes are derived from template.
    slide_data keys:
      - title (top centered)
      - body_left, body_right (string or list). If missing, split 'bullets' evenly.
    """
    # Layout ratios: 45% (left) + 10% (gap) + 45% (right)
    # No overlap by construction
    top_margin = Inches(0.2)
    bottom_margin = Inches(0.2)

    # Determine content
    title_text = slide_data.get("title", "")
    body_left = slide_data.get("body_left")
    body_right = slide_data.get("body_right")
    if body_left is None and body_right is None:
        bullets = slide_data.get("bullets", [])
        half = (len(bullets) + 1) // 2
        body_left = "\n".join(bullets[:half])
        body_right = "\n".join(bullets[half:])

    # Detect contrast and text color once
    contrast_type = get_contrast_text_color(slide, slide_width, slide_height)
    text_color = get_text_color_for_contrast(contrast_type, slide_data.get("theme", ["#003087", "#FFFFFF"]))

    # Scale fonts with slide width, but respect template inheritance if available
    scale_w = slide_width / Inches(10)  # 10in wide as baseline
    try:
        tmp = slide.shapes.add_textbox(Inches(0.1), Inches(0.1), Inches(1), Inches(0.5))
        tf = tmp.text_frame
        tf.text = "Tmp"
        inherited = get_inherited_font_size(tmp, slide, presentation)
        base_pt = inherited.pt if inherited else 20
        # Title/body base sizes and scaled
        title_pt = int(max(20, min(44, base_pt * 1.3 * scale_w)))
        body_pt = int(max(12, min(28, base_pt * 0.9 * scale_w)))
        slide.shapes._spTree.remove(tmp._element)
    except Exception:
        title_pt = int(28 * scale_w)
        body_pt = int(18 * scale_w)

    # Title area box (top centered)
    title_height = int(slide_height * 0.16)
    title_box = slide.shapes.add_textbox(Inches(0), top_margin, slide_width, title_height)
    ttf = title_box.text_frame
    ttf.clear()
    ttf.word_wrap = True
    p_title = ttf.paragraphs[0]
    p_title.text = title_text
    p_title.alignment = PP_ALIGN.CENTER
    for run in p_title.runs:
        run.font.bold = True
        run.font.size = Pt(title_pt)
        run.font.color.rgb = text_color

    # Columns area
    top_of_columns = top_margin + title_height + Inches(0.05)
    column_height = slide_height - top_of_columns - bottom_margin
    left_x = int(slide_width * 0.0)
    left_w = int(slide_width * 0.45)
    right_x = int(slide_width * 0.55)
    right_w = int(slide_width * 0.45)

    def body_to_lines(body_text):
        if isinstance(body_text, list):
            return body_text
        if isinstance(body_text, str):
            return [ln for ln in body_text.splitlines() if ln.strip()]
        return []

    def add_body_column(left_pos, width_pos, body_text):
        tx = slide.shapes.add_textbox(left_pos, top_of_columns, width_pos, column_height)
        tf = tx.text_frame
        tf.clear()
        tf.word_wrap = True
        lines = body_to_lines(body_text)
        if not lines:
            return
        # First paragraph
        tf.text = lines[0]
        for run in tf.paragraphs[0].runs:
            run.font.size = Pt(body_pt)
            run.font.color.rgb = text_color
        # Remaining as paragraphs
        for ln in lines[1:]:
            p = tf.add_paragraph()
            p.text = ln
            p.level = 1
            for run in p.runs:
                run.font.size = Pt(body_pt)
                run.font.color.rgb = text_color
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)

    # Left column near left border
    add_body_column(left_x, left_w, body_left)
    # Right column near right border
    add_body_column(right_x, right_w, body_right)

def generate_gamma_style_ppt(slides_content, template_path=None, layout_index=None):
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
        for i in range(len(prs.slides)-1, -1, -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]
    else:
        prs = Presentation()
    # Use user-selected layout if provided
    if layout_index is not None and 0 <= layout_index < len(prs.slide_layouts):
        base_layout = prs.slide_layouts[layout_index]
    else:
        base_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    for slide_data in slides_content:
        image_path = None
        image_desc = slide_data.get("image_desc")
        if image_desc:
            try:
                image_path = fetch_image_from_unsplash(image_desc)
            except Exception as e:
                print(f"Image fetch failed for '{image_desc}': {e}")
        # Add slide with selected or default layout
        slide = prs.slides.add_slide(base_layout)
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        padding = Inches(0.5)
        img_width = int(slide_width * 0.4)
        img_height = int(slide_height * 0.6)
        # Find first textbox placeholder (for text)
        txBox = None
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type in [1, 2, 14]:  # TITLE, BODY, CONTENT
                txBox = shape
                break
        
        if not txBox:
            txBox = slide.shapes.add_textbox(padding, padding, slide_width - 2*padding, slide_height - 2*padding)
        # If layout requests two-column text, render that and skip default single-column text
        layout_name = str(slide_data.get("layout", "")).strip().lower()
        if layout_name in ("two-column", "two column", "two columns"):
            add_two_column_content(slide, slide_data, output_ppt, content_slide_width, content_slide_height)
        else:
            tf = txBox.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.text = slide_data.get("title", "")
            for bullet in slide_data.get("bullets", []):
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 1
                p.font.size = Pt(18)
        
        # Image placement (if any) - Add image first so we can detect background
        if image_path:
            img_placeholder = None
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.type == 18:  # PICTURE
                    img_placeholder = shape
                    break
        if img_placeholder and os.path.exists(image_path):
            img_placeholder.insert_picture(image_path)
        else:
            if image_path and os.path.exists(image_path):
                img = slide.shapes.add_picture(image_path, padding, slide_height - img_height - padding, width=img_width, height=img_height)
                from PIL import Image as PILImage
                with PILImage.open(image_path) as pil_img:
                    aspect = pil_img.width / pil_img.height
                    if img_width / img_height > aspect:
                        new_width = int(img_height * aspect)
                        img.width = new_width
                        img.left = padding
                    else:
                        new_height = int(img_width / aspect)
                        img.height = new_height
                        img.top = slide_height - new_height - padding
        
        # Detect contrast and get appropriate text color AFTER adding image
        contrast_type = get_contrast_text_color(slide, slide_width, slide_height)
        text_color = get_text_color_for_contrast(contrast_type, slide_data.get("theme", ["#003087", "#FFFFFF"]))

        # Apply contrast-based text color to default single-column case
        if layout_name not in ("two-column", "two column", "two columns"):
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = text_color
            tf.margin_left = Inches(0.2)
            tf.margin_right = Inches(0.2)
            tf.margin_top = Inches(0.1)
            tf.margin_bottom = Inches(0.1)
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            for paragraph in tf.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
    output_path = os.path.join("outputs", f"gamma_style_{uuid.uuid4().hex[:8]}.pptx")
    prs.save(output_path)
    return output_path

def normalize_text(text):
    """Normalize text by removing extra spaces, newlines, and converting to lowercase for comparison."""
    if not text:
        return ""
    text = re.sub(r'\s+', ' ', text.strip()).lower()
    return text

def analyze_text_role(text, paragraph_level):
    """Analyze the role of text (e.g., title, body, bullet) using NLP and paragraph level."""
    doc = nlp(text) if 'nlp' in globals() and nlp else None
    if paragraph_level > 0:
        return "bullet"
    if any(line.strip().startswith(("-", "•", "*")) for line in text.splitlines()):
        return "bullet"
    if len(text.split()) < 5 or (doc is not None and len(doc.ents) > 0):
        return "title"
    else:
        return "body"

def get_inherited_font_size(shape, slide, presentation):
    """Get the inherited font size from shape, slide, or presentation theme if not explicitly set."""
    if shape.text_frame.paragraphs:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size is not None:
                    return run.font.size
    slide_layout = slide.slide_layout
    if slide_layout:
        for shape in slide_layout.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size is not None:
                            return run.font.size
    return Pt(18)

def get_inherited_font_color(shape, slide, presentation):
    """Get the inherited font color from shape, slide, or presentation theme if not explicitly set."""
    if shape.text_frame.paragraphs:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if hasattr(run.font.color, 'rgb') and run.font.color.rgb is not None:
                    return run.font.color.rgb
    slide_layout = slide.slide_layout
    if slide_layout:
        for shape in slide_layout.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if hasattr(run.font.color, 'rgb') and run.font.color.rgb is not None:
                            return run.font.color.rgb
    return RGBColor(0, 0, 0)

def copy_run_formatting(src_run, dst_run, shape, slide, presentation):
    # NO scaling, just copy original
    font_size = src_run.font.size
    if font_size is None:
        font_size = get_inherited_font_size(shape, slide, presentation)
    dst_run.font.size = font_size
    if src_run.font.name:
        dst_run.font.name = src_run.font.name
    if src_run.font.bold is not None:
        dst_run.font.bold = src_run.font.bold
    if src_run.font.italic is not None:
        dst_run.font.italic = src_run.font.italic
    if hasattr(src_run.font.color, 'rgb') and src_run.font.color.rgb is not None:
        dst_run.font.color.rgb = src_run.font.color.rgb
    else:
        dst_run.font.color.rgb = get_inherited_font_color(shape, slide, presentation)

def copy_paragraph_formatting(src_paragraph, dst_paragraph):
    if src_paragraph.alignment is not None:
        dst_paragraph.alignment = src_paragraph.alignment
    else:
        dst_paragraph.alignment = PP_ALIGN.LEFT
    if src_paragraph.space_before is not None:
        dst_paragraph.space_before = src_paragraph.space_before
    if src_paragraph.space_after is not None:
        dst_paragraph.space_after = src_paragraph.space_after
    if src_paragraph.line_spacing is not None:
        dst_paragraph.line_spacing = src_paragraph.line_spacing
    if src_paragraph.level is not None:
        dst_paragraph.level = src_paragraph.level

def copy_text_frame_formatting(src_text_frame, dst_text_frame, shape, slide, presentation):
    dst_text_frame.clear()
    if src_text_frame.margin_left is not None:
        dst_text_frame.margin_left = src_text_frame.margin_left
    if src_text_frame.margin_right is not None:
        dst_text_frame.margin_right = src_text_frame.margin_right
    if src_text_frame.margin_top is not None:
        dst_text_frame.margin_top = src_text_frame.margin_top
    if src_text_frame.margin_bottom is not None:
        dst_text_frame.margin_bottom = src_text_frame.margin_bottom
    if src_text_frame.vertical_anchor is not None:
        dst_text_frame.vertical_anchor = src_text_frame.vertical_anchor
    dst_text_frame.word_wrap = src_text_frame.word_wrap if src_text_frame.word_wrap is not None else True

    for src_paragraph in src_text_frame.paragraphs:
        dst_paragraph = dst_text_frame.add_paragraph()
        copy_paragraph_formatting(src_paragraph, dst_paragraph)
        for src_run in src_paragraph.runs:
            if src_run.text:
                dst_run = dst_paragraph.add_run()
                dst_run.text = src_run.text
                copy_run_formatting(src_run, dst_run, shape, slide, presentation)

def add_text_box(slide, src_shape, slide_width, slide_height, scale_factor, presentation, image_positions=None):
    text_box = slide.shapes.add_textbox(
        left=src_shape.left,
        top=src_shape.top,
        width=src_shape.width,
        height=src_shape.height
    )
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    copy_text_frame_formatting(src_shape.text_frame, text_frame, src_shape, slide, presentation)

def add_image_to_slide(slide, image_path, left, top, width, height, slide_width, slide_height, scale_factor):
    slide.shapes.add_picture(image_path, int(left), int(top), width=int(width), height=int(height))

def take_screenshot_of_master(ppt_path, output_folder):
    try:
        ppt_path = os.path.abspath(ppt_path)
        print(f"[DEBUG] Opening PPT file for screenshot: {ppt_path}")

        if not os.path.exists(ppt_path):
            raise FileNotFoundError(f"PPT template not found: {ppt_path}")

        # Ensure output folder exists
        os.makedirs(output_folder, exist_ok=True)

        # Windows path: use COM automation
        if WIN32_AVAILABLE:
            import pythoncom  # type: ignore
            import win32com.client  # type: ignore
            pythoncom.CoInitialize()
            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            ppt_app.Visible = True
            presentation = ppt_app.Presentations.Open(ppt_path, WithWindow=False)

            if presentation.Slides.Count < 1:
                raise Exception("No slides found in template PPT")
            slide = presentation.Slides(1)

            temp_image_path = os.path.join(output_folder, f"temp_{uuid.uuid4().hex[:8]}.png")
            temp_image_path = os.path.abspath(temp_image_path)
            slide.Export(temp_image_path, "PNG", 1280, 720)

            presentation.Close()
            ppt_app.Quit()
            pythoncom.CoUninitialize()

            if not os.path.exists(temp_image_path):
                raise Exception("Failed to export slide as image")
            print(f"[DEBUG] Slide screenshot exported (COM): {temp_image_path}")
            return temp_image_path

        # Non-Windows path: try LibreOffice if available
        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if soffice:
            try:
                # Kill any existing LibreOffice processes to avoid conflicts
                try:
                    subprocess.run(["pkill", "-f", "soffice"], capture_output=True, timeout=5)
                    import time
                    time.sleep(1)  # Wait for process to terminate
                except:
                    pass
                
                # Convert all slides to PNGs, then pick the first
                cmd = [soffice, "--headless", "--convert-to", "png", "--outdir", output_folder, ppt_path]
                print(f"[DEBUG] Running LibreOffice export: {' '.join(cmd)}")
                
                # Run with timeout and better error handling
                result = subprocess.run(
                    cmd, 
                    capture_output=True, 
                    text=True, 
                    timeout=30,
                    env=dict(os.environ, HOME=os.path.expanduser("~"))
                )
                
                if result.returncode != 0:
                    print(f"[DEBUG] LibreOffice stderr: {result.stderr}")
                    print(f"[DEBUG] LibreOffice stdout: {result.stdout}")
                
                base = os.path.splitext(os.path.basename(ppt_path))[0]
                # Find a generated PNG (pick the first)
                candidates = [f for f in os.listdir(output_folder) if f.startswith(base) and f.lower().endswith('.png')]
                if candidates:
                    first_png = os.path.join(output_folder, sorted(candidates)[0])
                    print(f"[DEBUG] Slide screenshot exported (LibreOffice): {first_png}")
                    return first_png
                else:
                    print(f"[DEBUG] No PNG files found in {output_folder}")
                    print(f"[DEBUG] Files in directory: {os.listdir(output_folder)}")
            except subprocess.TimeoutExpired:
                print("[DEBUG] LibreOffice export timed out")
            except Exception as e:
                print(f"LibreOffice export failed: {e}")
                import traceback
                traceback.print_exc()

        # Fallback: no screenshot capability
        print("[WARN] No screenshot method available on this platform. Proceeding without template background.")
        print("[INFO] To enable LibreOffice background processing:")
        print("[INFO] 1. Go to System Preferences > Security & Privacy > Privacy")
        print("[INFO] 2. Select 'Full Disk Access' and add LibreOffice")
        print("[INFO] 3. Select 'Automation' and ensure LibreOffice has permissions")
        return None
    except Exception as e:
        print(f"Error taking template slide screenshot: {e}")
        try:
            if 'presentation' in locals():
                presentation.Close()
            if 'ppt_app' in locals():
                ppt_app.Quit()
            if WIN32_AVAILABLE:
                import pythoncom  # type: ignore
                pythoncom.CoUninitialize()
        except:
            pass
        return None

def generate_ppt(content_path, template_path, layout_index=1):
    try:
        print(f"Template path: {template_path}")
        print(f"Content path: {content_path}")
        print(f"Layout index: {layout_index}")

        if not os.path.exists(template_path):
            raise FileNotFoundError("Template file not found")
        if content_path and not os.path.exists(content_path):
            raise FileNotFoundError("Content file not found")

        screenshot_dir = os.path.join("Uploads", "screenshots")
        os.makedirs(screenshot_dir, exist_ok=True)
        screenshot_path = take_screenshot_of_master(template_path, screenshot_dir)
        if not screenshot_path:
            print("[WARN] Failed to capture master slide screenshot. Continuing with template layouts (no screenshot background).")

        content_ppt = Presentation(content_path)
        template_ppt = Presentation(template_path)
        content_slide_width = content_ppt.slide_width
        content_slide_height = content_ppt.slide_height
        # If we couldn't take a screenshot, we still want to leverage the template's master/layouts
        if not screenshot_path:
            output_ppt = Presentation(template_path)
            # Remove any existing slides from the template so we add fresh ones with the same style
            for i in range(len(output_ppt.slides)-1, -1, -1):
                rId = output_ppt.slides._sldIdLst[i].rId
                output_ppt.part.drop_rel(rId)
                del output_ppt.slides._sldIdLst[i]
            print("[INFO] Using template background and layouts (no screenshot)")
        else:
            output_ppt = Presentation()
            print("[INFO] Using screenshot background")
        output_ppt.slide_width = content_slide_width
        output_ppt.slide_height = content_slide_height

        # Choose base layout: use provided layout_index if valid, else default
        try:
            base_idx = int(layout_index)
        except Exception:
            base_idx = None
        default_layout_index = 6 if len(output_ppt.slide_layouts) > 6 else 0
        if base_idx is not None and 0 <= base_idx < len(output_ppt.slide_layouts):
            base_layout = output_ppt.slide_layouts[base_idx]
        else:
            base_layout = output_ppt.slide_layouts[default_layout_index]
        print(f"Content PPT loaded with {len(content_ppt.slides)} slides")
        print(f"Content slide dimensions: {content_slide_width}x{content_slide_height}")
        print(f"Using layout index: {list(output_ppt.slide_layouts).index(base_layout)}")

        for slide_idx, slide in enumerate(content_ppt.slides):
            print(f"Processing slide {slide_idx + 1}")
            new_slide = output_ppt.slides.add_slide(base_layout)
            if screenshot_path and os.path.exists(screenshot_path):
                new_slide.shapes.add_picture(
                    screenshot_path, 
                    left=0, 
                    top=0, 
                    width=content_slide_width, 
                    height=content_slide_height
                )
                print(f"Added background image to slide {slide_idx + 1}")
            # If no screenshot, we rely on the template's background from the chosen layout
            # Remove default placeholders to avoid duplicates with copied content
            try:
                for shp in list(new_slide.shapes):
                    if getattr(shp, 'is_placeholder', False):
                        new_slide.shapes._spTree.remove(shp._element)
            except Exception:
                pass

            image_positions = []
            for shape in slide.shapes:
                if shape.shape_type == 13:
                    image_positions.append({
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height
                    })

            image_count = 0
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    add_text_box(new_slide, shape, content_slide_width, content_slide_height, 1.0, output_ppt, image_positions=image_positions)
                if shape.shape_type == 13:
                    image_count += 1
                    image_stream = shape.image.blob
                    image_filename = f"temp_image_{uuid.uuid4().hex[:8]}.png"
                    image_path = os.path.join("Uploads", "images", image_filename)
                    os.makedirs(os.path.dirname(image_path), exist_ok=True)
                    with open(image_path, 'wb') as f:
                        f.write(image_stream)
                    add_image_to_slide(
                        new_slide, image_path, 
                        shape.left, shape.top, 
                        shape.width, shape.height, 
                        content_slide_width, content_slide_height, 1.0
                    )
                    print(f"Added image to slide at adjusted position ({shape.left}, {shape.top})")

            print(f"Total images found in slide {slide_idx + 1}: {image_count}")

        output_dir = "outputs"
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, f"generated_{uuid.uuid4().hex[:8]}.pptx")
        output_ppt.save(output_path)
        print(f"PPT saved to: {output_path}")
        return output_path

    except Exception as e:
        print(f"Error generating PPT: {e}")
        raise Exception(f"Failed to generate PPT: {str(e)}")
    finally:
        try:
            if 'screenshot_path' in locals() and isinstance(screenshot_path, str) and screenshot_path and os.path.exists(screenshot_path):
                os.remove(screenshot_path)
        except Exception as _cleanup_err:
            print(f"[WARN] Could not remove temp screenshot: {_cleanup_err}")

# ...rest of your code (extract_ppt_content, clean_template_ppt, refine_ppt, etc.) remains unchanged...
def extract_ppt_content(ppt_path):
    """Extract content from a PPT file as structured data."""
    ppt = Presentation(ppt_path)
    content = []
    for slide_idx, slide in enumerate(ppt.slides):
        slide_content = {"slide": slide_idx + 1, "shapes": []}
        for shape in slide.shapes:
            shape_data = {}
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                shape_data["type"] = "text"
                shape_data["text"] = shape.text
                shape_data["position"] = {
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height
                }
                shape_data["formatting"] = []
                for paragraph in shape.text_frame.paragraphs:
                    para_data = {
                        "alignment": paragraph.alignment.value if paragraph.alignment else None,
                        "level": paragraph.level if paragraph.level is not None else 0
                    }
                    runs = []
                    for run in paragraph.runs:
                        run_data = {
                            "text": run.text,
                            "bold": run.font.bold if run.font.bold is not None else False,
                            "italic": run.font.italic if run.font.italic is not None else False,
                            "size": run.font.size.pt if run.font.size else None,
                            "color": str(run.font.color.rgb) if hasattr(run.font.color, 'rgb') and run.font.color.rgb else None
                        }
                        runs.append(run_data)
                    para_data["runs"] = runs
                    shape_data["formatting"].append(para_data)
            elif shape.shape_type == 13:
                shape_data["type"] = "image"
                shape_data["position"] = {
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height
                }
            if shape_data:
                slide_content["shapes"].append(shape_data)
        content.append(slide_content)
    return content

def clean_template_ppt(template_ppt):
    """Clean duplicate slide layouts, slide masters, themes, and media in the template PPT."""
    seen_layouts = set()
    seen_masters = set()
    seen_themes = set()
    seen_media = set()

    slide_layouts = list(template_ppt.slide_layouts)
    layout_ids_to_remove = []
    for layout in slide_layouts:
        layout_name = layout.name
        if layout_name in seen_layouts:
            layout_ids_to_remove.append(layout)
        else:
            seen_layouts.add(layout_name)
    
    for layout in layout_ids_to_remove:
        for r_id in list(template_ppt.part.rels.keys()):
            if template_ppt.part.rels[r_id].target_part == layout.part:
                template_ppt.part.drop_rel(r_id)
                break
        for sld_layout_id in template_ppt.slide_master.slide_layouts._sldLayoutIdLst:
            if sld_layout_id.rId == template_ppt.part.rels.get_id_of_part(layout.part):
                template_ppt.slide_master.slide_layouts._sldLayoutIdLst.remove(sld_layout_id)
                break
        if layout.part in template_ppt.part.related_parts.values():
            layout_part = layout.part
            template_ppt.part.package.drop_part(layout_part)

    slide_masters = list(template_ppt.slide_masters)
    master_ids_to_remove = []
    for master in slide_masters:
        master_name = master.name if hasattr(master, 'name') else str(id(master))
        if master_name in seen_masters:
            master_ids_to_remove.append(master)
        else:
            seen_masters.add(master_name)

    for master in master_ids_to_remove:
        for r_id in list(template_ppt.part.rels.keys()):
            if template_ppt.part.rels[r_id].target_part == master.part:
                template_ppt.part.drop_rel(r_id)
                break
        for sld_master_id in template_ppt.slide_masters._sldMasterIdLst:
            if sld_master_id.rId == template_ppt.part.rels.get_id_of_part(master.part):
                template_ppt.slide_masters._sldMasterIdLst.remove(sld_master_id)
                break
        if master.part in template_ppt.part.related_parts.values():
            master_part = master.part
            template_ppt.part.package.drop_part(master_part)

    themes = []
    for rel in template_ppt.part.rels.values():
        if rel.target_part.partname.startswith('/ppt/theme/theme'):
            themes.append(rel.target_part)
    
    for theme in themes[1:]:
        theme_name = theme.partname
        if theme_name in seen_themes:
            for r_id in list(template_ppt.part.rels.keys()):
                if template_ppt.part.rels[r_id].target_part == theme:
                    template_ppt.part.drop_rel(r_id)
                    break
            if theme in template_ppt.part.related_parts.values():
                template_ppt.part.package.drop_part(theme)
        else:
            seen_themes.add(theme_name)

    for slide in template_ppt.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:
                image_part = shape.image.part
                image_name = image_part.partname
                if image_name in seen_media:
                    new_image_name = f"/ppt/media/image_{uuid.uuid4().hex[:8]}.{image_name.split('.')[-1]}"
                    image_part.partname = new_image_name
                seen_media.add(image_part.partname)


def fill_template_with_ai_content(template_path, slides_content):
    prs = Presentation(template_path)
    layout_index = 6 if len(prs.slide_layouts) > 6 else 0
    blank_layout = prs.slide_layouts[layout_index]
    # Remove existing slides
    for i in range(len(prs.slides)-1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
    # Add AI slides
    for slide in slides_content:
        s = prs.slides.add_slide(blank_layout)
        txBox = s.shapes.add_textbox(Pt(100), Pt(100), Pt(700), Pt(400))
        tf = txBox.text_frame
        tf.text = slide['title']
        for bullet in slide['bullets']:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 1
    output_path = os.path.join("outputs", f"ai_generated_{uuid.uuid4().hex[:8]}.pptx")
    prs.save(output_path)
    return output_path

def fix_text_contrast_in_ppt(ppt_path):
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches
    import os

    def get_bg_brightness(slide, slide_width, slide_height):
        # Check for large image as background (assume light)
        for shape in slide.shapes:
            if shape.shape_type == 13:
                if (shape.left < Inches(0.5) and shape.top < Inches(0.5) and 
                    shape.width > slide_width * 0.8 and shape.height > slide_height * 0.8):
                    return 255  # Assume white bg for most templates
        # Check background fill safely (guard _NoFill)
        try:
            if hasattr(slide, 'background') and hasattr(slide.background, 'fill') and slide.background.fill:
                fill = slide.background.fill
                if hasattr(fill, 'fore_color'):
                    try:
                        fc = fill.fore_color
                        if getattr(fc, 'rgb', None):
                            r, g, b = fc.rgb
                            brightness = (r + g + b) / 3
                            return brightness
                    except Exception:
                        pass
        except Exception:
            pass
        return 255  # Default: white

    def get_text_brightness(rgb):
        if rgb is None:
            return 0
        r, g, b = rgb
        return (r + g + b) / 3

    prs = Presentation(ppt_path)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    changed = False
    for slide in prs.slides:
        bg_brightness = get_bg_brightness(slide, slide_width, slide_height)
        for shape in slide.shapes:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame and shape.text_frame.text.strip():
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        # Only adjust explicit colors
                        if getattr(run, 'font', None) and getattr(run.font, 'color', None) and getattr(run.font.color, 'rgb', None):
                            color = run.font.color.rgb
                            text_brightness = get_text_brightness(color)
                            if bg_brightness > 200 and text_brightness > 180:
                                run.font.color.rgb = RGBColor(0, 0, 0)
                                changed = True
                            elif bg_brightness < 80 and text_brightness < 100:
                                run.font.color.rgb = RGBColor(255, 255, 255)
                                changed = True
    if changed:
        out_path = os.path.join("outputs", f"contrast_fixed_{os.path.basename(ppt_path)}")
        prs.save(out_path)
        return out_path
    else:
        return ppt_path

def refine_ppt(converted_path, content_path, api_url, api_key):
    import shutil
    import os
    # Copy converted ppt as is
    temp_path = os.path.join("outputs", f"temp_{os.path.basename(converted_path)}")
    shutil.copyfile(converted_path, temp_path)
    # Fix color contrast
    out_path = fix_text_contrast_in_ppt(temp_path)
    return out_path

def generate_enhanced_ppt(slides_content, template_path=None, layout_preference="auto"):
    """
    Enhanced PPT generation with better layout handling and content fitting
    """
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
        # Clear existing slides but keep template structure
        for i in range(len(prs.slides)-1, -1, -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]
    else:
        prs = Presentation()
    
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    for i, slide_data in enumerate(slides_content):
        # Determine layout based on preference and slide type
        layout_type = determine_layout_type(slide_data, layout_preference, i)
        
        # Create slide with appropriate layout
        slide = create_slide_with_layout(prs, layout_type, slide_data, slide_width, slide_height)
        
        # Add content with proper fitting
        add_content_to_slide(slide, slide_data, layout_type, slide_width, slide_height)
    
    # Save the presentation
    output_filename = f"enhanced_ai_pptx_{uuid.uuid4().hex[:8]}.pptx"
    output_path = os.path.join("outputs", output_filename)
    os.makedirs("outputs", exist_ok=True)
    prs.save(output_path)
    
    return output_path

def determine_layout_type(slide_data, layout_preference, slide_index):
    """
    Determine the best layout type for a slide
    Prefers image layouts when an image description is provided
    """
    if layout_preference == "auto":
        title = slide_data.get("title", "").lower()
        bullets = slide_data.get("bullets", [])
        has_image_hint = bool(slide_data.get("image_desc"))

        if slide_index == 0 or "introduction" in title or "overview" in title:
            return "title_slide"

        # Prefer image layouts if image hint present
        if has_image_hint:
            # With moderate bullets, use side image; with very few, use full image
            if len(bullets) <= 5:
                return "image_left"
            else:
                return "image_right"

        # Fallbacks without image hints
        if len(bullets) <= 3:
            return "image_left"
        if len(bullets) > 5:
            return "two_column"
        return "title_content"

    layout_mapping = {
        "title-content": "title_content",
        "image-left": "image_left",
        "image-right": "image_right",
        "full-image": "full_image",
        "two-column": "two_column",
    }
    return layout_mapping.get(layout_preference, "title_content")

def create_slide_with_layout(prs, layout_type, slide_data, slide_width, slide_height):
    """
    Create a slide with the specified layout type
    """
    # Always use blank layout for better control over content placement
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    
    slide = prs.slides.add_slide(layout)
    return slide

def add_content_to_slide(slide, slide_data, layout_type, slide_width, slide_height):
    """
    Add content to slide with proper fitting based on layout type
    """
    # Completely clear the slide first
    clear_slide_completely(slide)
    
    title = slide_data.get("title", "")
    bullets = slide_data.get("bullets", [])
    image_desc = slide_data.get("image_desc", "")
    theme = slide_data.get("theme", ["#003087", "#FFFFFF"])
    
    # Fetch image if any hint present (even if some layouts might not use it)
    image_path = None
    if image_desc:
        try:
            image_path = fetch_image_from_unsplash(image_desc)
        except Exception as e:
            print(f"Image fetch failed for '{image_desc}': {e}")
    
    if layout_type == "title_slide":
        add_title_slide_content(slide, title, image_path, theme, slide_width, slide_height)
    elif layout_type == "title_content":
        add_title_content_layout(slide, title, bullets, theme, slide_width, slide_height)
    elif layout_type == "image_left":
        add_image_left_layout(slide, title, bullets, image_path, theme, slide_width, slide_height)
    elif layout_type == "image_right":
        add_image_right_layout(slide, title, bullets, image_path, theme, slide_width, slide_height)
    elif layout_type == "full_image":
        add_full_image_layout(slide, title, bullets, image_path, theme, slide_width, slide_height)
    elif layout_type == "two_column":
        add_two_column_layout(slide, title, bullets, theme, slide_width, slide_height)

def add_title_slide_content(slide, title, image_path, theme, slide_width, slide_height):
    """
    Add content to title slide
    """
    # Add background image if available
    if image_path:
        slide.shapes.add_picture(image_path, 0, 0, slide_width, slide_height)
        # Add semi-transparent overlay for text readability
        overlay = slide.shapes.add_shape(1, 0, 0, slide_width, slide_height)  # Rectangle
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
        overlay.fill.transparency = 0.3
        overlay.line.fill.background()
    
    # Detect contrast and get appropriate text color
    contrast_type = get_contrast_text_color(slide, slide_width, slide_height)
    text_color = get_text_color_for_contrast(contrast_type, theme)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_run = title_para.runs[0]
    title_run.font.size = Pt(44)
    title_run.font.bold = True
    title_run.font.color.rgb = text_color

def add_title_content_layout(slide, title, bullets, theme, slide_width, slide_height):
    """
    Add content to title and content layout
    """
    # Detect contrast and get appropriate text color
    contrast_type = get_contrast_text_color(slide, slide_width, slide_height)
    text_color = get_text_color_for_contrast(contrast_type, theme)
    
    # Compute scale relative to a typical 10in x 7.5in slide
    width_scale = max(0.7, min(1.5, slide_width / Inches(10)))
    height_scale = max(0.7, min(1.5, slide_height / Inches(7.5)))
    scale = (width_scale + height_scale) / 2

    # Heuristic body size adjustment based on bullet count
    bullet_factor = 1.0
    if bullets:
        # More bullets -> smaller size, but keep within [0.75, 1.0]
        bullet_factor = max(0.75, min(1.0, 6 / max(3, len(bullets))))

    # Dynamic font sizes
    title_font_size = Pt(int(34 * scale))
    body_font_size = Pt(int(19 * scale * bullet_factor))

    # Add title: full-width box, centered at top, bold
    title_margin_h = Inches(0.7)
    title_box_width = slide_width - 2 * title_margin_h
    title_box_left = title_margin_h
    title_box_top = Inches(0.5)
    title_box_height = Inches(1.6)
    title_box = slide.shapes.add_textbox(title_box_left, title_box_top, title_box_width, title_box_height)
    title_frame = title_box.text_frame
    title_frame.clear()
    title_frame.word_wrap = True
    title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_run = title_para.runs[0]
    title_run.font.size = title_font_size
    title_run.font.bold = True
    title_run.font.color.rgb = text_color
    
    # Add bullets: centered in middle area, non-bold
    content_margin_h = Inches(1.0)
    content_width = slide_width - 2 * content_margin_h
    content_height = int(slide_height * 0.45)
    content_left = content_margin_h
    # Vertically center in remaining space (below title)
    remaining_top = title_box_top + title_box_height
    remaining_height = slide_height - remaining_top - Inches(0.6)
    content_top = remaining_top + max(0, (remaining_height - content_height) // 2)
    bullets_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    bullets_frame = bullets_box.text_frame
    bullets_frame.clear()
    bullets_frame.word_wrap = True
    
    for idx, bullet in enumerate(bullets or []):
        p = bullets_frame.add_paragraph() if idx > 0 else bullets_frame.paragraphs[0]
        p.text = bullet
        p.level = 0
        p.alignment = PP_ALIGN.CENTER
        # spacing for readability
        p.space_after = Pt(6 * scale)
        for run in p.runs:
            run.font.size = body_font_size
            run.font.bold = False
            run.font.color.rgb = text_color

def add_image_left_layout(slide, title, bullets, image_path, theme, slide_width, slide_height):
    """
    Add content with image on left, text on right
    """
    # Scale relative to baseline 10in x 7.5in
    width_scale = max(0.7, min(1.5, slide_width / Inches(10)))
    height_scale = max(0.7, min(1.5, slide_height / Inches(7.5)))
    scale = (width_scale + height_scale) / 2

    # Heuristic body size adjustment based on bullet count
    bullet_factor = 1.0
    if bullets:
        bullet_factor = max(0.75, min(1.0, 6 / max(3, len(bullets))))

    # Dynamic font sizes
    title_font_size = Pt(int(28 * scale))
    body_font_size = Pt(int(18 * scale * bullet_factor))

    # Detect contrast and get appropriate text color
    contrast_type = get_contrast_text_color(slide, slide_width, slide_height)
    text_color = get_text_color_for_contrast(contrast_type, theme)

    # Layout metrics
    left_margin = Inches(0.0)  # image flush with left border
    right_margin = Inches(0.4)
    gap_between = Inches(0.35)
    top_margin = Inches(0.4)
    bottom_margin = Inches(0.4)

    # Add image on left (maintain aspect ratio, fill height or width as needed)
    img_box_left = left_margin
    img_box_top = top_margin
    img_box_height_max = slide_height - top_margin - bottom_margin
    img_box_width_max = int(slide_width * 0.44)  # ~44% width to leave room for text

    img_rendered_width = None
    img_rendered_height = None
    if image_path and os.path.exists(image_path):
        try:
            with Image.open(image_path) as im:
                aspect = im.width / im.height if im.height else 1.0
        except Exception:
            aspect = 16/9

        # Try width-constrained first
        w1 = img_box_width_max
        h1 = int(w1 / aspect)
        if h1 > img_box_height_max:
            # Height-constrained
            h1 = int(img_box_height_max)
            w1 = int(h1 * aspect)
        img_rendered_width = w1
        img_rendered_height = h1

        # Center vertically within the box
        img_top = img_box_top + max(0, (img_box_height_max - img_rendered_height) // 2)
        img_left = img_box_left  # flush to left border
        slide.shapes.add_picture(image_path, int(img_left), int(img_top), width=int(img_rendered_width), height=int(img_rendered_height))

    # Right content region
    content_left = (img_box_left + (img_rendered_width or 0) + gap_between) if image_path else Inches(1)
    content_top = top_margin
    content_width = slide_width - content_left - right_margin
    content_height = slide_height - top_margin - bottom_margin

    # Title on top of right pane
    title_height = max(Inches(0.9), Inches(1.2) * scale)
    title_box = slide.shapes.add_textbox(int(content_left), int(content_top), int(content_width), int(title_height))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_frame.word_wrap = True
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_run = title_para.runs[0]
    title_para.alignment = PP_ALIGN.LEFT
    title_run.font.size = title_font_size
    title_run.font.bold = True
    title_run.font.color.rgb = text_color

    # Bullets in remaining space of right pane
    bullets_top = content_top + title_height + Inches(0.1)
    bullets_height = content_height - title_height - Inches(0.1)
    bullets_box = slide.shapes.add_textbox(int(content_left), int(bullets_top), int(content_width), int(bullets_height))
    bullets_frame = bullets_box.text_frame
    bullets_frame.clear()
    bullets_frame.word_wrap = True
    
    for idx, bullet in enumerate(bullets or []):
        p = bullets_frame.add_paragraph() if idx > 0 else bullets_frame.paragraphs[0]
        p.text = bullet
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.size = body_font_size
            run.font.bold = False
            run.font.color.rgb = text_color

def add_image_right_layout(slide, title, bullets, image_path, theme, slide_width, slide_height):
    """
    Add content with image on right, text on left
    """
    # Scale relative to baseline 10in x 7.5in
    width_scale = max(0.7, min(1.5, slide_width / Inches(10)))
    height_scale = max(0.7, min(1.5, slide_height / Inches(7.5)))
    scale = (width_scale + height_scale) / 2

    # Heuristic body size adjustment based on bullet count
    bullet_factor = 1.0
    if bullets:
        bullet_factor = max(0.75, min(1.0, 6 / max(3, len(bullets))))

    # Dynamic font sizes
    title_font_size = Pt(int(28 * scale))
    body_font_size = Pt(int(18 * scale * bullet_factor))

    # Detect contrast and get appropriate text color
    contrast_type = get_contrast_text_color(slide, slide_width, slide_height)
    text_color = get_text_color_for_contrast(contrast_type, theme)

    # Layout metrics
    left_margin = Inches(0.4)
    right_margin = Inches(0.0)  # image flush with right border
    gap_between = Inches(0.35)
    top_margin = Inches(0.4)
    bottom_margin = Inches(0.4)

    # Image on right: compute rendered size and flush to right border
    img_box_height_max = slide_height - top_margin - bottom_margin
    img_box_width_max = int(slide_width * 0.44)

    img_rendered_width = None
    img_rendered_height = None
    if image_path and os.path.exists(image_path):
        try:
            with Image.open(image_path) as im:
                aspect = im.width / im.height if im.height else 1.0
        except Exception:
            aspect = 16/9

        w1 = img_box_width_max
        h1 = int(w1 / aspect)
        if h1 > img_box_height_max:
            h1 = int(img_box_height_max)
            w1 = int(h1 * aspect)
        img_rendered_width = w1
        img_rendered_height = h1

        img_top = top_margin + max(0, (img_box_height_max - img_rendered_height) // 2)
        # flush to right
        img_left = slide_width - right_margin - img_rendered_width
        slide.shapes.add_picture(image_path, int(img_left), int(img_top), width=int(img_rendered_width), height=int(img_rendered_height))

    # Left content region
    content_left = left_margin
    content_top = top_margin
    content_width = slide_width - left_margin - right_margin - (img_rendered_width or 0) - gap_between
    content_height = slide_height - top_margin - bottom_margin

    # Title on top of left pane
    title_height = max(Inches(0.9), Inches(1.2) * scale)
    title_box = slide.shapes.add_textbox(int(content_left), int(content_top), int(content_width), int(title_height))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_frame.word_wrap = True
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_run = title_para.runs[0]
    title_para.alignment = PP_ALIGN.LEFT
    title_run.font.size = title_font_size
    title_run.font.bold = True
    title_run.font.color.rgb = text_color

    # Bullets in remaining space of left pane
    bullets_top = content_top + title_height + Inches(0.1)
    bullets_height = content_height - title_height - Inches(0.1)
    bullets_box = slide.shapes.add_textbox(int(content_left), int(bullets_top), int(content_width), int(bullets_height))
    bullets_frame = bullets_box.text_frame
    bullets_frame.clear()
    bullets_frame.word_wrap = True

    for idx, bullet in enumerate(bullets or []):
        p = bullets_frame.add_paragraph() if idx > 0 else bullets_frame.paragraphs[0]
        p.text = bullet
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.size = body_font_size
            run.font.bold = False
            run.font.color.rgb = text_color

def add_full_image_layout(slide, title, bullets, image_path, theme, slide_width, slide_height):
    """
    Add content with full image background
    """
    # Background image full-bleed
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, 0, 0, slide_width, slide_height)

    # Dynamic scaling based on slide size and bullet count
    width_scale = max(0.7, min(1.5, slide_width / Inches(10)))
    height_scale = max(0.7, min(1.5, slide_height / Inches(7.5)))
    scale = (width_scale + height_scale) / 2
    bullet_factor = 1.0
    if bullets:
        bullet_factor = max(0.75, min(1.0, 6 / max(3, len(bullets))))

    # Font sizes with sensible clamps
    def clamp(v, lo, hi):
        return max(lo, min(hi, v))
    title_pt = clamp(int(36 * scale), 18, 60)
    body_pt = clamp(int(20 * scale * bullet_factor), 12, 30)
    title_font_size = Pt(title_pt)
    body_font_size = Pt(body_pt)

    # Bottom-centered translucent black strip (container for title + bullets)
    strip_left = int(slide_width * 0.08)
    strip_width = int(slide_width * 0.84)
    strip_height = int(slide_height * 0.32)
    strip_top = int(slide_height - strip_height - int(slide_height * 0.06))

    text_bg = slide.shapes.add_shape(1, strip_left, strip_top, strip_width, strip_height)
    text_bg.fill.solid()
    text_bg.fill.fore_color.rgb = RGBColor(0, 0, 0)
    text_bg.fill.transparency = 0.7
    text_bg.line.fill.background()

    # White text for contrast on black strip
    text_color = RGBColor(255, 255, 255)

    # Title: top-center inside strip, bold (fixed region height + gap to bullets)
    title_pad_h = int(strip_height * 0.18)
    title_box = slide.shapes.add_textbox(
        strip_left + int(strip_width * 0.05),
        strip_top + int(strip_height * 0.08),
        int(strip_width * 0.90),
        title_pad_h
    )
    title_frame = title_box.text_frame
    title_frame.clear()
    title_frame.word_wrap = True
    title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    title_frame.margin_left = Pt(4)
    title_frame.margin_right = Pt(4)
    title_frame.margin_top = Pt(2)
    title_frame.margin_bottom = Pt(2)
    title_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_run = title_para.runs[0]
    title_run.font.size = title_font_size
    title_run.font.bold = True
    title_run.font.color.rgb = text_color

    # Bullets: centered within remaining strip area (with fixed gap from title)
    gap_between_title_bullets = int(strip_height * 0.08)
    bullets_box_top = strip_top + int(strip_height * 0.08) + title_pad_h + gap_between_title_bullets
    bullets_box_height = strip_top + strip_height - bullets_box_top - int(strip_height * 0.06)
    bullets_box = slide.shapes.add_textbox(
        strip_left + int(strip_width * 0.08),
        bullets_box_top,
        int(strip_width * 0.84),
        bullets_box_height
    )
    bullets_frame = bullets_box.text_frame
    bullets_frame.clear()
    bullets_frame.word_wrap = True
    bullets_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    bullets_frame.margin_left = Pt(4)
    bullets_frame.margin_right = Pt(4)
    bullets_frame.margin_top = Pt(2)
    bullets_frame.margin_bottom = Pt(2)
    bullets_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    for idx, bullet in enumerate(bullets or []):
        p = bullets_frame.add_paragraph() if idx > 0 else bullets_frame.paragraphs[0]
        p.text = bullet
        p.level = 0
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(6 * scale)
        for run in p.runs:
            run.font.size = body_font_size
            run.font.bold = False
            run.font.color.rgb = text_color

def add_two_column_layout(slide, title, bullets, theme, slide_width, slide_height):
    """
    Add content in two-column layout
    """
    # Detect contrast and get appropriate text color
    contrast_type = get_contrast_text_color(slide, slide_width, slide_height)
    text_color = get_text_color_for_contrast(contrast_type, theme)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_run = title_para.runs[0]
    title_run.font.size = Pt(32)
    title_run.font.bold = True
    title_run.font.color.rgb = text_color
    
    # Split bullets into two columns
    mid_point = len(bullets) // 2
    left_bullets = bullets[:mid_point]
    right_bullets = bullets[mid_point:]
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3.5), Inches(5))
    left_frame = left_box.text_frame
    
    for bullet in left_bullets:
        p = left_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(16)
            run.font.color.rgb = text_color
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(3.5), Inches(5))
    right_frame = right_box.text_frame
    
    for bullet in right_bullets:
        p = right_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(16)
            run.font.color.rgb = text_color

def get_contrast_text_color(slide, slide_width, slide_height):
    """
    Detect slide background color and return appropriate text color for contrast
    Returns: 'light' for dark backgrounds, 'dark' for light backgrounds
    """
    try:
        # First check if there's a background image
        has_background_image = any(shape.shape_type == 13 for shape in slide.shapes)
        
        # Check if slide has a background fill
        background = slide.background
        if hasattr(background, 'fill') and background.fill.type == 1:  # Solid fill
            fill_color = background.fill.fore_color
            if hasattr(fill_color, 'rgb'):
                rgb = fill_color.rgb
                if rgb:
                    # Simple brightness calculation (0-255)
                    brightness = (rgb[0] * 299 + rgb[1] * 587 + rgb[2] * 114) / 1000
                    # If background is light, use dark text; if dark, use light text
                    return 'dark' if brightness > 150 else 'light'
        
        # If there's a background image, default to light text for better visibility
        if has_background_image:
            return 'light'
            
        # Check slide layout background if no specific background is set
        if hasattr(slide, 'slide_layout'):
            layout = slide.slide_layout
            if hasattr(layout, 'background') and hasattr(layout.background, 'fill') and layout.background.fill.type == 1:
                fill_color = layout.background.fill.fore_color
                if hasattr(fill_color, 'rgb'):
                    rgb = fill_color.rgb
                    if rgb:
                        brightness = (rgb[0] * 299 + rgb[1] * 587 + rgb[2] * 114) / 1000
                        return 'dark' if brightness > 150 else 'light'
        
        # Default: assume light background (white) - return dark text
        return 'dark'
        
    except Exception as e:
        print(f"Error detecting background color: {e}")
        # Default fallback - assume light background, return dark text
        return 'dark'

def get_text_color_for_contrast(contrast_type, theme_colors=None):
    """
    Get appropriate text color based on contrast type
    
    Args:
        contrast_type: 'light' for dark backgrounds, 'dark' for light backgrounds
        theme_colors: Optional list of theme colors (not used in this simplified version)
    
    Returns:
        RGBColor object with appropriate text color
    """
    if contrast_type == 'light':
        # For dark backgrounds, use white text
        return RGBColor(255, 255, 255)  # White
    else:
        # For light backgrounds, use black text for maximum contrast
        return RGBColor(0, 0, 0)  # Black

def clear_slide_completely(slide):
    """
    Completely clear all shapes and placeholders from a slide
    """
    try:
        # Remove all shapes from the slide
        shapes_to_remove = []
        for shape in slide.shapes:
            shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            try:
                slide.shapes._spTree.remove(shape._element)
            except:
                pass
        
        # Clear background
        if hasattr(slide, 'background'):
            try:
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
            except:
                pass
                
    except Exception as e:
        print(f"Error clearing slide: {e}")
